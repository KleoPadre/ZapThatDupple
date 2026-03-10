from pathlib import Path
from typing import Optional

DOCUMENT_EXTENSIONS = {
    ".pdf", ".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx",
    ".txt", ".md", ".rtf", ".odt", ".ods", ".odp", ".csv",
    ".html", ".htm", ".xml", ".json", ".yaml", ".yml",
}

MAX_TEXT_CHARS = 10_000  # limit for embedding


def extract_text(path: str) -> Optional[str]:
    ext = Path(path).suffix.lower()

    try:
        if ext == ".pdf":
            return _extract_pdf(path)
        elif ext in (".doc", ".docx"):
            return _extract_docx(path)
        elif ext in (".xls", ".xlsx"):
            return _extract_excel(path)
        elif ext in (".ppt", ".pptx"):
            return _extract_pptx(path)
        elif ext == ".rtf":
            return _extract_rtf(path)
        elif ext in (".txt", ".md", ".csv", ".html", ".htm", ".xml", ".json", ".yaml", ".yml"):
            return _extract_plain(path)
        else:
            return _extract_plain(path)
    except Exception as e:
        print(f"Text extraction error {path}: {e}")
        return None


def _extract_pdf(path: str) -> str:
    import fitz  # pymupdf
    doc = fitz.open(path)
    texts = []
    for page in doc:
        texts.append(page.get_text())
    return " ".join(texts)[:MAX_TEXT_CHARS]


def _extract_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    return " ".join(p.text for p in doc.paragraphs)[:MAX_TEXT_CHARS]


def _extract_excel(path: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    texts = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    texts.append(str(cell))
    return " ".join(texts)[:MAX_TEXT_CHARS]


def _extract_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return " ".join(texts)[:MAX_TEXT_CHARS]


def _extract_rtf(path: str) -> str:
    with open(path, "rb") as f:
        content = f.read()
    # Basic RTF strip
    import re
    text = content.decode("utf-8", errors="ignore")
    text = re.sub(r"\{[^{}]*\}", "", text)
    text = re.sub(r"\\[a-z]+\d*\s?", " ", text)
    return text[:MAX_TEXT_CHARS]


def _extract_plain(path: str) -> str:
    import chardet
    with open(path, "rb") as f:
        raw = f.read(MAX_TEXT_CHARS * 4)
    detected = chardet.detect(raw)
    encoding = detected.get("encoding") or "utf-8"
    return raw.decode(encoding, errors="ignore")[:MAX_TEXT_CHARS]
